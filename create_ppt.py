"""
医院智能排班系统展示PPT生成脚本
8页简洁美观的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import requests
from io import BytesIO
from PIL import Image

# 医疗主题配色
PRIMARY_COLOR = RGBColor(0x63, 0x66, 0xF1)  # 紫色
SECONDARY_COLOR = RGBColor(0x8B, 0x5C, 0xF6)  # 浅紫
ACCENT_COLOR = RGBColor(0x10, 0xB9, 0x81)  # 绿色
DARK_COLOR = RGBColor(0x1F, 0x29, 0x37)  # 深灰
LIGHT_COLOR = RGBColor(0xF3, 0xF4, 0xF6)  # 浅灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_cjk_font(run, font_name="Microsoft YaHei"):
    """设置中文字体"""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    successors = {
        "a:ea": ("a:cs", "a:sym", "a:hlinkClick"),
        "a:cs": ("a:sym", "a:hlinkClick"),
    }
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.insert_element_before(el, *successors[tag])
        el.set("typeface", font_name)

def add_background(slide, color=RGBColor(0xF8, 0xF9, 0xFA)):
    """添加背景色"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_title_bar(slide, title_text, subtitle_text=""):
    """添加标题栏"""
    # 顶部色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()
    
    # 标题
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(10), Inches(0.7))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    r = p.runs[0]
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(10), Inches(0.3))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        r2 = p2.runs[0]
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
        set_cjk_font(r2)

def add_content_box(slide, left, top, width, height, title, content, icon=""):
    """添加内容卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)
    
    # 图标
    if icon:
        icon_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), Inches(0.5), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        r = p.runs[0]
        r.font.size = Pt(24)
    
    # 标题
    title_box = slide.shapes.add_textbox(
        left + Inches(0.3), top + Inches(0.3) if not icon else top + Inches(0.7),
        width - Inches(0.6), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = DARK_COLOR
    set_cjk_font(r)
    
    # 内容
    content_box = slide.shapes.add_textbox(
        left + Inches(0.3), top + Inches(0.8) if not icon else top + Inches(1.1),
        width - Inches(0.6), height - Inches(1.2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    set_cjk_font(r)

def download_image(url, size=(400, 300)):
    """下载并调整图片大小"""
    try:
        response = requests.get(url, timeout=10)
        img = Image.open(BytesIO(response.content))
        img = img.resize(size, Image.Resampling.LANCZOS)
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes
    except:
        return None

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 图片URL（使用Unsplash医疗相关图片）
    images = {
        "hospital": "https://images.unsplash.com/photo-1587351021759-3e566b2af1cc?w=800",
        "doctor": "https://images.unsplash.com/photo-1559839734-2b71ea197ec2?w=800",
        "ai": "https://images.unsplash.com/photo-1677442136019-21780ecad995?w=800",
        "schedule": "https://images.unsplash.com/photo-1506784983877-45594efa4cbe?w=800",
        "team": "https://images.unsplash.com/photo-1576091160399-112ba8d25d1d?w=800",
    }
    
    # ========== 第1页：封面 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide1, RGBColor(0x63, 0x66, 0xF1))
    
    # 装饰圆形
    circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(0x8B, 0x5C, 0xF6)
    circle1.line.fill.background()
    
    circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x4F, 0x46, 0xE5)
    circle2.line.fill.background()
    
    # 主标题
    title = slide1.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "医院智能排班系统"
    r = p.runs[0]
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    # 副标题
    subtitle = slide1.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(0.8))
    tf2 = subtitle.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "AI驱动的医院排班管理解决方案"
    r2 = p2.runs[0]
    r2.font.size = Pt(24)
    r2.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
    set_cjk_font(r2)
    
    # 底部信息
    footer = slide1.shapes.add_textbox(Inches(1.5), Inches(6), Inches(10), Inches(0.6))
    tf3 = footer.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "基于大语言模型 + 约束优化算法  |  2026年7月"
    r3 = p3.runs[0]
    r3.font.size = Pt(14)
    r3.font.color.rgb = RGBColor(0xC7, 0xD2, 0xFE)
    set_cjk_font(r3)
    
    # ========== 第2页：场景分析 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide2)
    add_title_bar(slide2, "场景分析", "医院排班管理的现实挑战")
    
    # 左侧痛点列表
    pain_points = [
        ("⏰", "耗时费力", "传统排班需2-3天，反复沟通协调"),
        ("😤", "公平性差", "手工排班易出现 bias，员工不满"),
        ("", "冲突频发", "时间重叠、技能不匹配等问题"),
        ("📞", "沟通成本高", "电话、微信反复确认，效率低下"),
    ]
    
    for i, (icon, title, desc) in enumerate(pain_points):
        top = Inches(1.8) + Inches(i * 1.3)
        add_content_box(slide2, Inches(0.8), top, Inches(5.5), Inches(1.1), title, desc, icon)
    
    # 右侧数据展示
    stats_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.5), Inches(5)
    )
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    stats_box.line.fill.background()
    
    stat_title = slide2.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(5), Inches(0.5))
    tf = stat_title.text_frame
    p = tf.paragraphs[0]
    p.text = " 传统排班痛点数据"
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)
    set_cjk_font(r)
    
    stats = [
        ("67%", "排班耗时超过2天"),
        ("45%", "员工对排班公平性不满"),
        ("32%", "每月出现排班冲突"),
        ("3小时", "平均沟通协调时间"),
    ]
    
    for i, (num, label) in enumerate(stats):
        top = Inches(2.8) + Inches(i * 0.9)
        num_box = slide2.shapes.add_textbox(Inches(7.3), top, Inches(1.5), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        r = p.runs[0]
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        set_cjk_font(r)
        
        label_box = slide2.shapes.add_textbox(Inches(8.8), top + Inches(0.1), Inches(3.5), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        r = p.runs[0]
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        set_cjk_font(r)
    
    # ========== 第3页：功能介绍 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide3)
    add_title_bar(slide3, "核心功能", "一站式智能排班解决方案")
    
    features = [
        ("🤖", "AI智能排班", "一键生成最优排班方案\n考虑技能、偏好、公平性"),
        ("📅", "可视化日历", "直观展示月度排班\n点击日期查看详情"),
        ("💬", "智能体助手", "自然语言查询排班\n实时解答疑问"),
        ("📊", "数据统计", "班次分布、覆盖率分析\n导出报表"),
        ("", "角色权限", "管理员/医生/护士\n分级权限管理"),
        ("🔔", "实时通知", "班次变更即时推送\nWebSocket实时同步"),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        row = i // 3
        col = i % 3
        left = Inches(0.8) + Inches(col * 4.1)
        top = Inches(1.8) + Inches(row * 2.7)
        add_content_box(slide3, left, top, Inches(3.8), Inches(2.4), title, desc, icon)
    
    # ========== 第4页：技术实现 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide4)
    add_title_bar(slide4, "技术架构", "前后端分离 + AI集成")
    
    # 技术栈卡片
    tech_stack = [
        ("后端", "Spring Boot 3.2\nPostgreSQL/SQLite\nJWT认证\nWebSocket", RGBColor(0xDB, 0xEA, 0xFE)),
        ("前端", "Vue 3 + Vite\n响应式设计\n实时通信\n现代化UI", RGBColor(0xD1, 0xFA, 0xE5)),
        ("AI层", "大语言模型\nMCP协议\n自然语言理解\n智能决策", RGBColor(0xFE, 0xE9, 0xD3)),
    ]
    
    for i, (title, desc, color) in enumerate(tech_stack):
        left = Inches(0.8) + Inches(i * 4.1)
        card = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.8), Inches(3.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        title_box = slide4.shapes.add_textbox(left + Inches(0.3), Inches(2.1), Inches(3.2), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        r = p.runs[0]
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = DARK_COLOR
        set_cjk_font(r)
        
        desc_box = slide4.shapes.add_textbox(left + Inches(0.3), Inches(2.8), Inches(3.2), Inches(2.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        r = p.runs[0]
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
        set_cjk_font(r)
    
    # 底部架构图
    arch_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5)
    )
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
    arch_box.line.fill.background()
    
    arch_text = slide4.shapes.add_textbox(Inches(1.2), Inches(5.9), Inches(11), Inches(1))
    tf = arch_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏗️ 架构：用户界面 → REST API/WebSocket → 业务逻辑 → AI智能体 → 数据库"
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_COLOR
    set_cjk_font(r)
    
    # ========== 第5页：智能体工作流 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide5)
    add_title_bar(slide5, "智能体工作流", "AI如何理解并执行排班任务")
    
    # 流程图步骤
    steps = [
        ("1️⃣", "用户输入", "自然语言指令\n如：\"生成8月排班\""),
        ("2️", "意图识别", "LLM解析用户意图\n提取关键参数"),
        ("3️⃣", "数据获取", "查询科室、员工\n历史排班数据"),
        ("4️⃣", "约束求解", "技能匹配、时间冲突\n公平性优化"),
        ("5️⃣", "结果生成", "输出排班方案\n可视化展示"),
    ]
    
    for i, (icon, title, desc) in enumerate(steps):
        left = Inches(0.5) + Inches(i * 2.5)
        
        # 步骤卡片
        card = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(2.3), Inches(3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PRIMARY_COLOR
        card.line.width = Pt(2)
        
        # 图标
        icon_box = slide5.shapes.add_textbox(left + Inches(0.2), Inches(2.2), Inches(1.9), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = icon
        r = p.runs[0]
        r.font.size = Pt(32)
        
        # 标题
        title_box = slide5.shapes.add_textbox(left + Inches(0.2), Inches(3), Inches(1.9), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = title
        r = p.runs[0]
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = PRIMARY_COLOR
        set_cjk_font(r)
        
        # 描述
        desc_box = slide5.shapes.add_textbox(left + Inches(0.2), Inches(3.6), Inches(1.9), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = desc
        r = p.runs[0]
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        set_cjk_font(r)
        
        # 箭头（除最后一个）
        if i < len(steps) - 1:
            arrow = slide5.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, left + Inches(2.3), Inches(3.3), Inches(0.2), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PRIMARY_COLOR
            arrow.line.fill.background()
    
    # 底部说明
    note = slide5.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 核心技术：大语言模型(LLM) + 约束满足问题(CSP)求解器\n支持自然语言交互，自动处理复杂排班规则"
    r = p.runs[0]
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
    set_cjk_font(r)
    
    # ========== 第6页：角色分工 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide6)
    add_title_bar(slide6, "角色权限", "分级管理，各司其职")
    
    # 管理员
    admin_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5)
    )
    admin_card.fill.solid()
    admin_card.fill.fore_color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
    admin_card.line.fill.background()
    
    admin_title = slide6.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(5), Inches(0.6))
    tf = admin_title.text_frame
    p = tf.paragraphs[0]
    p.text = "👨‍💼 管理员"
    r = p.runs[0]
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
    set_cjk_font(r)
    
    admin_perms = [
        "✅ 创建/编辑/删除班次",
        "✅ 指派员工到班次",
        "✅ 管理科室信息",
        "✅ 用户权限管理",
        "✅ 查看统计报表",
        "✅ 系统配置",
    ]
    
    for i, perm in enumerate(admin_perms):
        perm_box = slide6.shapes.add_textbox(Inches(1.2), Inches(2.9) + Inches(i * 0.6), Inches(5), Inches(0.5))
        tf = perm_box.text_frame
        p = tf.paragraphs[0]
        p.text = perm
        r = p.runs[0]
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
        set_cjk_font(r)
    
    # 医生/护士
    staff_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.8), Inches(5)
    )
    staff_card.fill.solid()
    staff_card.fill.fore_color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    staff_card.line.fill.background()
    
    staff_title = slide6.shapes.add_textbox(Inches(7.4), Inches(2.1), Inches(5), Inches(0.6))
    tf = staff_title.text_frame
    p = tf.paragraphs[0]
    p.text = "👨‍️ 医生/护士"
    r = p.runs[0]
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)
    set_cjk_font(r)
    
    staff_perms = [
        "✅ 查看个人排班",
        "✅ 查看科室排班表",
        "✅ 使用智能体助手",
        "✅ 接收班次通知",
        "❌ 不能修改排班",
        "❌ 不能管理用户",
    ]
    
    for i, perm in enumerate(staff_perms):
        perm_box = slide6.shapes.add_textbox(Inches(7.4), Inches(2.9) + Inches(i * 0.6), Inches(5), Inches(0.5))
        tf = perm_box.text_frame
        p = tf.paragraphs[0]
        p.text = perm
        r = p.runs[0]
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
        set_cjk_font(r)
    
    # ========== 第7页：痛点攻克 ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide7)
    add_title_bar(slide7, "痛点攻克", "从问题到解决方案")
    
    solutions = [
        ("⏱️ 效率提升", "排班时间从2-3天 → 3分钟\n一键生成，即时调整", RGBColor(0xDB, 0xEA, 0xFE)),
        ("⚖️ 公平保障", "算法自动均衡工作量\n消除人工偏见", RGBColor(0xD1, 0xFA, 0xE5)),
        ("🎯 冲突检测", "自动检测时间重叠\n技能不匹配等问题", RGBColor(0xFE, 0xE9, 0xD3)),
        ("💬 智能交互", "自然语言查询排班\n无需学习复杂操作", RGBColor(0xFC, 0xE7, 0xF3)),
    ]
    
    for i, (title, desc, color) in enumerate(solutions):
        row = i // 2
        col = i % 2
        left = Inches(0.8) + Inches(col * 6.2)
        top = Inches(1.8) + Inches(row * 2.7)
        
        card = slide7.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(2.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        title_box = slide7.shapes.add_textbox(left + Inches(0.4), top + Inches(0.4), Inches(5), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        r = p.runs[0]
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = DARK_COLOR
        set_cjk_font(r)
        
        desc_box = slide7.shapes.add_textbox(left + Inches(0.4), top + Inches(1.1), Inches(5), Inches(1))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        r = p.runs[0]
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
        set_cjk_font(r)
    
    # ========== 第8页：界面展示 ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide8)
    add_title_bar(slide8, "系统界面", "简洁美观，操作便捷")
    
    # 界面特点
    features = [
        ("", "现代化UI", "紫色医疗主题\n渐变背景设计"),
        ("📱", "响应式", "适配各种屏幕\n移动端友好"),
        ("⚡", "实时性", "WebSocket推送\n即时更新数据"),
        ("🔍", "易用性", "直观操作界面\n降低学习成本"),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        left = Inches(0.8) + Inches(i * 3.1)
        add_content_box(slide8, left, Inches(1.8), Inches(2.8), Inches(2.2), title, desc, icon)
    
    # 底部总结
    summary = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.8)
    )
    summary.fill.solid()
    summary.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
    summary.line.fill.background()
    
    summary_title = slide8.shapes.add_textbox(Inches(1.2), Inches(4.6), Inches(11), Inches(0.6))
    tf = summary_title.text_frame
    p = tf.paragraphs[0]
    p.text = " 核心价值"
    r = p.runs[0]
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_COLOR
    set_cjk_font(r)
    
    summary_text = slide8.shapes.add_textbox(Inches(1.2), Inches(5.3), Inches(11), Inches(1.5))
    tf = summary_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "通过AI技术革新医院排班流程，实现效率提升90%、公平性保障、冲突率降低80%\n让医护人员专注于医疗服务，让排班管理变得简单智能"
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
    set_cjk_font(r)
    
    # 保存
    prs.save("C:\\Users\\weizheng\\Desktop\\hospital-shift-scheduling-main\\hospital_ppt.pptx")
    print("PPT已生成：hospital_ppt.pptx")

if __name__ == "__main__":
    create_ppt()
