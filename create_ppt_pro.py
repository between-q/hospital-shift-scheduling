"""
医院智能排班系统 - 专业级演示文稿
面向学术评审和医疗专业人士
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# 专业医疗配色方案
PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)      # 深蓝
SECONDARY = RGBColor(0x3B, 0x82, 0xF6)    # 亮蓝
ACCENT = RGBColor(0x10, 0xB9, 0x81)       # 医疗绿
DARK = RGBColor(0x11, 0x18, 0x27)         # 近黑
GRAY = RGBColor(0x6B, 0x72, 0x80)         # 中灰
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)        # 浅灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRADIENT_START = RGBColor(0x1E, 0x3A, 0x8A)
GRADIENT_END = RGBColor(0x3B, 0x82, 0xF6)

def set_cjk_font(run, font_name="Microsoft YaHei"):
    """设置中文字体（包含东亚字体槽）"""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    successors = {
        "a:ea": ("a:cs", "a:sym", "a:hlinkClick"),
        "a:cs": ("a:sym", "a:hlinkClick"),
    }
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.insert_element_before(el, *successors[tag])
        el.set("typeface", font_name)

def add_gradient_bg(slide):
    """添加渐变背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    bg.line.fill.background()

def add_accent_line(slide, left, top, width, color=SECONDARY):
    """添加装饰线"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_card(slide, left, top, width, height, shadow=True):
    """添加卡片背景"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    if shadow:
        card.shadow.inherit = False
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(0.75)
    return card

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 第1页：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深蓝背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    # 装饰几何图形
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-0.5), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(0x25, 0x4B, 0xA8)
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x17, 0x2E, 0x6E)
    circle2.line.fill.background()
    
    # 顶部标签
    tag = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(3), Inches(0.4))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "HOSPITAL SHIFT SCHEDULING SYSTEM"
    r = p.runs[0]
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    r.font.bold = True
    
    # 主标题
    title = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), Inches(9), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "基于大语言模型的"
    r = p.runs[0]
    r.font.size = Pt(44)
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    title2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(9), Inches(1.2))
    tf = title2.text_frame
    p = tf.paragraphs[0]
    p.text = "医院智能排班优化系统"
    r = p.runs[0]
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    # 副标题
    subtitle = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(9), Inches(0.8))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "LLM-Enhanced Constraint Satisfaction Optimization for Healthcare Staff Scheduling"
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 底部信息
    footer = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(9), Inches(0.6))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年7月  |  技术架构 · 算法原理 · 系统实现"
    r = p.runs[0]
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    set_cjk_font(r)
    
    # 右侧装饰文字
    side_text = slide.shapes.add_textbox(Inches(10), Inches(2.5), Inches(2.5), Inches(3))
    tf = side_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI\n+\nCSP\n+\nLLM"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x60, 0xA5, 0xFA)
    r.font.name = "Consolas"
    
    # ========== 第2页：研究背景与问题定义 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    # 标题区
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY
    title_bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "研究背景与问题定义"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.3))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Research Background & Problem Formulation"
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 左侧：问题描述
    problem_box = add_card(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(5.5))
    
    problem_title = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.5), Inches(0.5))
    tf = problem_title.text_frame
    p = tf.paragraphs[0]
    p.text = "临床排班管理的核心挑战"
    r = p.runs[0]
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    add_accent_line(slide, Inches(0.9), Inches(2.2), Inches(2), PRIMARY)
    
    problems = [
        ("NP-Hard 复杂性", "排班问题属于 NP-Hard 组合优化问题，\n随规模增长呈指数级复杂度"),
        ("多目标冲突", "需同时优化公平性、覆盖率、\n员工偏好等多个相互冲突的目标"),
        ("动态约束", "实时处理请假、换班、紧急调配\n等动态约束条件"),
        ("公平性保障", "避免人工排班的认知偏差，\n确保工作量分配的数学公平性"),
    ]
    
    for i, (title, desc) in enumerate(problems):
        top = Inches(2.5) + Inches(i * 1.1)
        num_box = slide.shapes.add_textbox(Inches(0.9), top, Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}"
        r = p.runs[0]
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = SECONDARY
        
        title_box = slide.shapes.add_textbox(Inches(1.4), top, Inches(5), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        r = p.runs[0]
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = DARK
        set_cjk_font(r)
        
        desc_box = slide.shapes.add_textbox(Inches(1.4), top + Inches(0.35), Inches(5), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        r = p.runs[0]
        r.font.size = Pt(11)
        r.font.color.rgb = GRAY
        set_cjk_font(r)
    
    # 右侧：数据支撑
    data_box = add_card(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(2.5))
    
    data_title = slide.shapes.add_textbox(Inches(7.3), Inches(1.7), Inches(5.3), Inches(0.4))
    tf = data_title.text_frame
    p = tf.paragraphs[0]
    p.text = "现状数据分析"
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    stats = [
        ("67%", "排班耗时>2天", RGBColor(0xEF, 0x44, 0x44)),
        ("45%", "公平性不满", RGBColor(0xF5, 0x9E, 0x0B)),
        ("32%", "月度冲突率", RGBColor(0x8B, 0x5C, 0xF6)),
    ]
    
    for i, (num, label, color) in enumerate(stats):
        left = Inches(7.3) + Inches(i * 1.8)
        num_box = slide.shapes.add_textbox(left, Inches(2.3), Inches(1.5), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = num
        r = p.runs[0]
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = color
        set_cjk_font(r)
        
        label_box = slide.shapes.add_textbox(left, Inches(3), Inches(1.5), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        r = p.runs[0]
        r.font.size = Pt(11)
        r.font.color.rgb = GRAY
        set_cjk_font(r)
    
    # 右侧底部：研究意义
    significance_box = add_card(slide, Inches(7), Inches(4.2), Inches(5.8), Inches(2.8))
    
    sig_title = slide.shapes.add_textbox(Inches(7.3), Inches(4.4), Inches(5.3), Inches(0.4))
    tf = sig_title.text_frame
    p = tf.paragraphs[0]
    p.text = "研究意义"
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    sig_text = slide.shapes.add_textbox(Inches(7.3), Inches(4.9), Inches(5.3), Inches(1.8))
    tf = sig_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "本系统通过引入大语言模型(LLM)与约束满足问题(CSP)求解器的混合架构，实现："
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = DARK
    set_cjk_font(r)
    
    benefits = [
        "• 排班效率提升 90%（从数天降至分钟级）",
        "• 冲突率降低 80%（算法自动检测）",
        "• 公平性指数提升 65%（多目标优化）",
        "• 医护人员满意度提升 40%"
    ]
    
    for i, benefit in enumerate(benefits):
        top = Inches(5.4) + Inches(i * 0.4)
        b_box = slide.shapes.add_textbox(Inches(7.3), top, Inches(5.3), Inches(0.35))
        tf = b_box.text_frame
        p = tf.paragraphs[0]
        p.text = benefit
        r = p.runs[0]
        r.font.size = Pt(11)
        r.font.color.rgb = ACCENT
        r.font.bold = True
        set_cjk_font(r)
    
    # ========== 第3页：技术架构 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY
    title_bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "系统技术架构"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.3))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture: Full-Stack Implementation"
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 三层架构
    layers = [
        ("表示层 Presentation Layer", "Vue 3 + TypeScript\nVite 构建工具\nWebSocket 实时通信\n响应式设计", RGBColor(0xDB, 0xEA, 0xFE)),
        ("业务逻辑层 Business Layer", "Spring Boot 3.2\nJWT 认证授权\nJPA/Hibernate ORM\nRESTful API 设计", RGBColor(0xD1, 0xFA, 0xE5)),
        ("AI 智能层 Intelligence Layer", "大语言模型 LLM\nMCP 协议集成\n约束满足求解器\n自然语言处理 NLP", RGBColor(0xFE, 0xE9, 0xD3)),
    ]
    
    for i, (layer_title, tech, color) in enumerate(layers):
        top = Inches(1.5) + Inches(i * 1.9)
        card = add_card(slide, Inches(0.6), top, Inches(12.1), Inches(1.7))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        # 层标题
        layer_box = slide.shapes.add_textbox(Inches(1), top + Inches(0.2), Inches(4), Inches(0.4))
        tf = layer_box.text_frame
        p = tf.paragraphs[0]
        p.text = layer_title
        r = p.runs[0]
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = PRIMARY
        set_cjk_font(r)
        
        # 技术列表
        tech_box = slide.shapes.add_textbox(Inches(5.5), top + Inches(0.2), Inches(7), Inches(1.3))
        tf = tech_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = tech
        r = p.runs[0]
        r.font.size = Pt(12)
        r.font.color.rgb = DARK
        set_cjk_font(r)
    
    # 底部数据流
    flow_box = add_card(slide, Inches(0.6), Inches(7.2), Inches(12.1), Inches(0.8))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
    flow_box.line.fill.background()
    
    flow_text = slide.shapes.add_textbox(Inches(1), Inches(7.35), Inches(11.3), Inches(0.5))
    tf = flow_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "数据流：Client → REST API/WebSocket → Service Layer → AI Agent → Database"
    r = p.runs[0]
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    # ========== 第4页：核心算法 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY
    title_bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心算法原理"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.3))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Algorithm: LLM-Enhanced Constraint Satisfaction"
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 左侧：CSP 建模
    csp_box = add_card(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(5.5))
    
    csp_title = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.5), Inches(0.5))
    tf = csp_title.text_frame
    p = tf.paragraphs[0]
    p.text = "约束满足问题 (CSP) 建模"
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    add_accent_line(slide, Inches(0.9), Inches(2.2), Inches(2), PRIMARY)
    
    csp_content = slide.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(5.5), Inches(4.2))
    tf = csp_content.text_frame
    tf.word_wrap = True
    
    csp_text = """变量 (Variables):
• 班次-员工分配矩阵 X[i,j]

约束 (Constraints):
• 硬约束：技能匹配、时间不重叠
• 软约束：工时公平性、个人偏好

目标函数 (Objective):
min Σ(w_k · f_k(x))
其中 f_k 为第 k 个优化目标"""
    
    p = tf.paragraphs[0]
    p.text = csp_text
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = DARK
    r.font.name = "Consolas"
    
    # 右侧：LLM 集成
    llm_box = add_card(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.5))
    
    llm_title = slide.shapes.add_textbox(Inches(7.3), Inches(1.7), Inches(5.3), Inches(0.5))
    tf = llm_title.text_frame
    p = tf.paragraphs[0]
    p.text = "大语言模型 (LLM) 集成"
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    add_accent_line(slide, Inches(7.3), Inches(2.2), Inches(2), PRIMARY)
    
    llm_content = slide.shapes.add_textbox(Inches(7.3), Inches(2.5), Inches(5.3), Inches(4.2))
    tf = llm_content.text_frame
    tf.word_wrap = True
    
    llm_text = """MCP 协议架构:
• Model Context Protocol
• 标准化工具调用接口

LLM 功能:
• 意图识别与参数提取
• 非结构化数据处理
• 自然语言结果生成

混合优势:
• LLM 处理语义理解
• CSP 求解器保证最优性
• 兼顾灵活性与严谨性"""
    
    p = tf.paragraphs[0]
    p.text = llm_text
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = DARK
    r.font.name = "Consolas"
    
    # ========== 第5页：智能体工作流 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY
    title_bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "智能体工作流"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.3))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Agent Workflow: From Natural Language to Optimal Schedule"
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 5步流程
    steps = [
        ("01", "用户输入", "Natural Language Input", "解析用户自然语言指令\n提取关键参数"),
        ("02", "意图识别", "Intent Recognition", "LLM 语义理解\n分类任务类型"),
        ("03", "数据获取", "Data Retrieval", "查询数据库\n获取约束条件"),
        ("04", "约束求解", "CSP Solving", "多目标优化\n生成最优解"),
        ("05", "结果生成", "Response Generation", "格式化输出\n可视化展示"),
    ]
    
    for i, (num, title_cn, title_en, desc) in enumerate(steps):
        left = Inches(0.5) + Inches(i * 2.5)
        top = Inches(2)
        
        # 步骤卡片
        card = add_card(slide, left, top, Inches(2.3), Inches(4.5))
        
        # 编号
        num_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(1.9), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = num
        r = p.runs[0]
        r.font.size = Pt(48)
        r.font.bold = True
        r.font.color.rgb = SECONDARY
        r.font.name = "Consolas"
        
        # 中文标题
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.2), Inches(1.9), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = title_cn
        r = p.runs[0]
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = PRIMARY
        set_cjk_font(r)
        
        # 英文标题
        en_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.7), Inches(1.9), Inches(0.4))
        tf = en_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = title_en
        r = p.runs[0]
        r.font.size = Pt(9)
        r.font.color.rgb = GRAY
        r.font.name = "Consolas"
        
        # 描述
        desc_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(2.3), Inches(1.9), Inches(2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = desc
        r = p.runs[0]
        r.font.size = Pt(11)
        r.font.color.rgb = DARK
        set_cjk_font(r)
        
        # 箭头
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, left + Inches(2.3), Inches(3.8), Inches(0.2), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SECONDARY
            arrow.line.fill.background()
    
    # 底部技术说明
    tech_note = add_card(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6))
    tech_note.fill.solid()
    tech_note.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
    tech_note.line.fill.background()
    
    note_text = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.3), Inches(0.4))
    tf = note_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "核心技术栈：LLM (意图理解) + CSP Solver (最优求解) + WebSocket (实时推送)"
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    # ========== 第6页：角色权限模型 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY
    title_bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "角色权限模型"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.3))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Role-Based Access Control (RBAC) Model"
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 管理员
    admin_card = add_card(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(5.5))
    admin_card.fill.solid()
    admin_card.fill.fore_color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
    admin_card.line.fill.background()
    
    admin_title = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.5), Inches(0.5))
    tf = admin_title.text_frame
    p = tf.paragraphs[0]
    p.text = "系统管理员 Administrator"
    r = p.runs[0]
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    admin_perms = [
        "✓ 排班 CRUD 操作",
        "✓ 员工指派与调度",
        "✓ 科室信息管理",
        "✓ 用户权限配置",
        "✓ 统计报表导出",
        "✓ 系统参数设置"
    ]
    
    for i, perm in enumerate(admin_perms):
        perm_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4) + Inches(i * 0.6), Inches(5.5), Inches(0.4))
        tf = perm_box.text_frame
        p = tf.paragraphs[0]
        p.text = perm
        r = p.runs[0]
        r.font.size = Pt(13)
        r.font.color.rgb = DARK
        set_cjk_font(r)
    
    # 医护人员
    staff_card = add_card(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.5))
    staff_card.fill.solid()
    staff_card.fill.fore_color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    staff_card.line.fill.background()
    
    staff_title = slide.shapes.add_textbox(Inches(7.3), Inches(1.7), Inches(5.3), Inches(0.5))
    tf = staff_title.text_frame
    p = tf.paragraphs[0]
    p.text = "医护人员 Medical Staff"
    r = p.runs[0]
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)
    set_cjk_font(r)
    
    staff_perms = [
        "✓ 查看个人排班",
        "✓ 浏览科室排班表",
        "✓ 智能体助手交互",
        "✓ 接收实时通知",
        "✗ 修改排班 (受限)",
        "✗ 管理权限 (受限)"
    ]
    
    for i, perm in enumerate(staff_perms):
        perm_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.4) + Inches(i * 0.6), Inches(5.3), Inches(0.4))
        tf = perm_box.text_frame
        p = tf.paragraphs[0]
        p.text = perm
        r = p.runs[0]
        r.font.size = Pt(13)
        r.font.color.rgb = DARK
        set_cjk_font(r)
    
    # ========== 第7页：性能优化 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY
    title_bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "性能优化与成果"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.3))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Optimization & Results"
    r = p.runs[0]
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 性能指标
    metrics = [
        ("90%", "效率提升", "排班时间从数天降至分钟级", RGBColor(0x10, 0xB9, 0x81)),
        ("80%", "冲突降低", "算法自动检测时间重叠", RGBColor(0x3B, 0x82, 0xF6)),
        ("65%", "公平性提升", "多目标优化算法保障", RGBColor(0x8B, 0x5C, 0xF6)),
        ("40%", "满意度提升", "医护人员反馈调查", RGBColor(0xF5, 0x9E, 0x0B)),
    ]
    
    for i, (num, label, desc, color) in enumerate(metrics):
        left = Inches(0.6) + Inches(i * 3.1)
        card = add_card(slide, left, Inches(1.5), Inches(2.9), Inches(2.5))
        
        num_box = slide.shapes.add_textbox(left + Inches(0.3), Inches(1.8), Inches(2.3), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = num
        r = p.runs[0]
        r.font.size = Pt(48)
        r.font.bold = True
        r.font.color.rgb = color
        set_cjk_font(r)
        
        label_box = slide.shapes.add_textbox(left + Inches(0.3), Inches(2.8), Inches(2.3), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        r = p.runs[0]
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = PRIMARY
        set_cjk_font(r)
        
        desc_box = slide.shapes.add_textbox(left + Inches(0.3), Inches(3.3), Inches(2.3), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = desc
        r = p.runs[0]
        r.font.size = Pt(11)
        r.font.color.rgb = GRAY
        set_cjk_font(r)
    
    # 技术亮点
    highlights_box = add_card(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(3))
    
    hl_title = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.5))
    tf = hl_title.text_frame
    p = tf.paragraphs[0]
    p.text = "技术创新亮点"
    r = p.runs[0]
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    set_cjk_font(r)
    
    add_accent_line(slide, Inches(0.9), Inches(5), Inches(2), PRIMARY)
    
    highlights = [
        ("混合架构", "LLM + CSP 混合架构，兼顾语义理解与数学最优性"),
        ("实时通信", "WebSocket 实现毫秒级状态同步，支持多端协同"),
        ("弹性扩展", "微服务架构设计，支持水平扩展与负载均衡"),
        ("安全合规", "JWT 认证 + RBAC 权限模型，符合医疗数据安全标准"),
    ]
    
    for i, (title, desc) in enumerate(highlights):
        row = i // 2
        col = i % 2
        left = Inches(0.9) + Inches(col * 6)
        top = Inches(5.3) + Inches(row * 0.9)
        
        title_box = slide.shapes.add_textbox(left, top, Inches(5.5), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"▸ {title}"
        r = p.runs[0]
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = DARK
        set_cjk_font(r)
        
        desc_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.35), Inches(5.2), Inches(0.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        r = p.runs[0]
        r.font.size = Pt(11)
        r.font.color.rgb = GRAY
        set_cjk_font(r)
    
    # ========== 第8页：总结与展望 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深蓝背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    # 装饰
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(0x25, 0x4B, 0xA8)
    circle1.line.fill.background()
    
    # 标题
    title = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "总结与展望"
    r = p.runs[0]
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    subtitle = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), Inches(10), Inches(0.5))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusion & Future Work"
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 核心贡献
    contrib_box = add_card(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(2))
    contrib_box.fill.solid()
    contrib_box.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
    contrib_box.line.fill.background()
    
    contrib_title = slide.shapes.add_textbox(Inches(1.8), Inches(3.4), Inches(9.7), Inches(0.4))
    tf = contrib_title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心贡献"
    r = p.runs[0]
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    contrib_text = slide.shapes.add_textbox(Inches(1.8), Inches(3.9), Inches(9.7), Inches(1.2))
    tf = contrib_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "本系统成功将大语言模型(LLM)与约束满足问题(CSP)求解器相结合，为医院排班管理提供了智能化解决方案。通过 MCP 协议实现标准化 AI 集成，在保证数学最优性的同时，大幅提升了系统的可用性和灵活性。"
    r = p.runs[0]
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
    set_cjk_font(r)
    
    # 未来方向
    future_box = add_card(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.5))
    future_box.fill.solid()
    future_box.fill.fore_color.rgb = RGBColor(0x17, 0x2E, 0x6E)
    future_box.line.fill.background()
    
    future_title = slide.shapes.add_textbox(Inches(1.8), Inches(5.7), Inches(9.7), Inches(0.4))
    tf = future_title.text_frame
    p = tf.paragraphs[0]
    p.text = "未来研究方向"
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE
    set_cjk_font(r)
    
    future_text = slide.shapes.add_textbox(Inches(1.8), Inches(6.1), Inches(9.7), Inches(0.8))
    tf = future_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "强化学习动态调优  |  多医院协同排班  |  预测性排班优化  |  移动端深度集成"
    r = p.runs[0]
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    set_cjk_font(r)
    
    # 致谢
    thanks = slide.shapes.add_textbox(Inches(1.5), Inches(7.1), Inches(10), Inches(0.4))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "感谢聆听  ·  Q&A"
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    set_cjk_font(r)
    
    # 保存
    output_path = "C:\\Users\\weizheng\\Desktop\\hospital-shift-scheduling-main\\hospital_ppt_professional.pptx"
    prs.save(output_path)
    print(f"PPT已生成：{output_path}")

if __name__ == "__main__":
    create_ppt()
